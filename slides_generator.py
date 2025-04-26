from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import Dict, List

class SlidesGenerator:
    def __init__(self):
        """Initialize the PowerPoint presentation generator."""
        self.prs = Presentation()
        self._setup_slide_layouts()

    def _setup_slide_layouts(self):
        """Set up the slide layouts and store them as instance variables."""
        # Title slide layout
        self.title_slide_layout = self.prs.slide_layouts[0]
        # Content slide layout with title and bullet points
        self.content_slide_layout = self.prs.slide_layouts[1]

    def add_title_slide(self, title: str) -> None:
        """Add a title slide to the presentation."""
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        
        # Add subtitle placeholder for ads
        subtitle = slide.placeholders[1]
        subtitle.text = "[Advertisement Space]"

    def add_content_slide(self, title: str, points: List[str]) -> None:
        """Add a content slide with a title and bullet points."""
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        
        # Add and format title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        
        # Add and format bullet points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.level = 0

    def generate_presentation(self, content: Dict[str, List[str]], output_path: str) -> None:
        """Generate a complete PowerPoint presentation."""
        # Add title slide
        self.add_title_slide(content['title'][0])
        
        # Process content slides
        current_title = None
        current_points = []
        
        for item in content['content']:
            # If this item is in outline, it's a title
            if item in content['outline']:
                # If we have accumulated points, add the previous slide
                if current_title and current_points:
                    self.add_content_slide(current_title, current_points)
                # Start a new slide
                current_title = item
                current_points = []
            else:
                # Add to current points
                current_points.append(item)
        
        # Add the last slide if there are remaining points
        if current_title and current_points:
            self.add_content_slide(current_title, current_points)
        
        # Save the presentation
        self.prs.save(output_path)
