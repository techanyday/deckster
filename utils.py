import openai
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
import tempfile

def check_user_limits(user):
    """Check if user has reached their presentation limits."""
    if user.subscription_type == 'free':
        # Free users can create up to 3 presentations per day
        today_presentations = len([p for p in user.presentations if p.created_at.date() == datetime.utcnow().date()])
        return today_presentations < 3
    return True

def get_max_slides(user):
    """Get maximum number of slides based on user's subscription."""
    limits = {
        'free': 5,
        'pro': 15,
        'business': 30
    }
    return limits.get(user.subscription_type, 5)

def add_watermark(image_path, watermark_text="Created with Windsurf"):
    """Add watermark to an image."""
    try:
        # Open the image
        with Image.open(image_path) as img:
            # Create a copy to draw on
            watermarked = img.copy()
            draw = ImageDraw.Draw(watermarked)
            
            # Calculate text size and position
            width, height = img.size
            font_size = int(min(width, height) * 0.03)  # Scale font size with image
            try:
                font = ImageFont.truetype("arial.ttf", font_size)
            except:
                font = ImageFont.load_default()
                
            text_width = draw.textlength(watermark_text, font=font)
            text_height = font_size
            
            # Position text in bottom right corner with padding
            padding = 10
            x = width - text_width - padding
            y = height - text_height - padding
            
            # Add semi-transparent white background for text
            text_bbox = (x-padding, y-padding, x+text_width+padding, y+text_height+padding)
            draw.rectangle(text_bbox, fill=(255, 255, 255, 128))
            
            # Draw text
            draw.text((x, y), watermark_text, fill=(100, 100, 100), font=font)
            
            # Save the watermarked image
            watermarked.save(image_path)
            return True
    except Exception as e:
        print(f"Error adding watermark: {str(e)}")
        return False

def generate_presentation_content(input_text, mode='outline'):
    """Generate presentation content using OpenAI."""
    try:
        system_prompt = {
            'outline': """You are an expert presentation creator. Create a presentation outline from the given input. 
                      Format the output as a JSON array of slides. Each slide should have a 'title' and 'content'.
                      Keep titles concise and content focused. Include key points and avoid lengthy paragraphs.""",
            'bullet': """You are an expert at converting text into presentation bullet points. 
                     Create concise, impactful bullet points from the input text.
                     Format the output as a JSON array of slides. Each slide should have a 'title' and 'content'.
                     Keep bullet points short and focused."""
        }

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k",
            messages=[
                {"role": "system", "content": system_prompt.get(mode, system_prompt['outline'])},
                {"role": "user", "content": f"Create a presentation from this text: {input_text}"}
            ],
            temperature=0.7
        )
        
        return response.choices[0].message['content']
    except Exception as e:
        print(f"Error generating content: {str(e)}")
        return None

def create_ppt(content, template='modern'):
    """Create a PowerPoint presentation."""
    try:
        # Create presentation
        prs = Presentation()
        
        # Define slide layouts based on template
        layouts = {
            'modern': {
                'title_slide': prs.slide_layouts[0],
                'content_slide': prs.slide_layouts[1]
            }
        }
        
        current_template = layouts.get(template, layouts['modern'])
        
        # Process each slide
        for slide_data in content:
            # Create new slide
            slide = prs.slides.add_slide(current_template['content_slide'])
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data['title']
            
            # Set content
            content_shape = slide.placeholders[1]
            content_shape.text = slide_data['content']
            
            # Apply formatting
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.alignment = PP_ALIGN.LEFT
                
        # Save to temporary file
        temp_ppt = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_ppt.name)
        
        return temp_ppt
    except Exception as e:
        print(f"Error creating PowerPoint: {str(e)}")
        return None
