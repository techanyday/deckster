from flask import Flask, jsonify
from test_models import db, TestUser
from flask_login import LoginManager, login_required
from flask_cors import CORS
import openai
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///test.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['SECRET_KEY'] = 'test-secret-key'

# Set OpenAI API key
openai.api_key = os.getenv('OPENAI_API_KEY', 'sk-test123')

db.init_app(app)
CORS(app)

login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

@login_manager.user_loader
def load_user(user_id):
    return TestUser.query.get(int(user_id))

@app.route('/')
def hello():
    return 'Hello, World!'

@app.route('/protected')
@login_required
def protected():
    return 'Protected Area'

@app.route('/test_openai')
def test_openai():
    try:
        client = openai.Client()
        return jsonify({'status': 'OpenAI client initialized'})
    except Exception as e:
        return jsonify({'error': str(e)})

@app.route('/test_pptx')
def test_pptx():
    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Test Presentation"
        return jsonify({'status': 'PowerPoint slide created'})
    except Exception as e:
        return jsonify({'error': str(e)})

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True)
