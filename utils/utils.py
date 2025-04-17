import os
import io
from datetime import datetime
from models import SubscriptionPlan, Presentation
from openai import OpenAI
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import requests

def check_user_limits(user):
    # Get user's subscription plan
    plan = SubscriptionPlan.query.get(user.subscription_plan_id)
    if not plan:
        return False, "No subscription plan found"
    
    # Get count of presentations created this month
    start_of_month = datetime.utcnow().replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    presentations_this_month = Presentation.query.filter(
        Presentation.user_id == user.id,
        Presentation.created_at >= start_of_month
    ).count()
    
    # Check if user has exceeded their monthly limit
    if presentations_this_month >= plan.monthly_presentations:
        return False, f"Monthly limit of {plan.monthly_presentations} presentations reached"
    
    return True, None

def get_max_slides(user):
    plan = SubscriptionPlan.query.get(user.subscription_plan_id)
    return plan.slides_per_presentation if plan else 5

def add_watermark(image_path, watermark_text):
    # Open the image
    with Image.open(image_path) as img:
        # Create a copy to draw on
        watermarked = img.copy()
        
        # Create drawing context
        draw = ImageDraw.Draw(watermarked)
        
        # Calculate text size and position
        width, height = img.size
        try:
            font = ImageFont.truetype("arial.ttf", 30)
        except:
            font = ImageFont.load_default()
            
        # Add watermark text
        text_width = draw.textlength(watermark_text, font=font)
        x = width - text_width - 10
        y = height - 40
        
        # Draw text with shadow for better visibility
        draw.text((x+2, y+2), watermark_text, font=font, fill='black')
        draw.text((x, y), watermark_text, font=font, fill='white')
        
        # Save to bytes
        img_byte_arr = io.BytesIO()
        watermarked.save(img_byte_arr, format=img.format)
        img_byte_arr.seek(0)
        
        return img_byte_arr

def generate_presentation_content(topic, num_slides):
    client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    prompt = f"""Create a presentation outline for {topic} with {num_slides} slides.
    For each slide include:
    - A title
    - 3-4 key points
    - Any relevant image suggestions
    Format as a JSON array where each object has: title, points (array), and image_suggestion."""
    
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a presentation expert. Create clear, engaging slides."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error generating content: {str(e)}")
        return None

def create_ppt(content, watermark=None):
    prs = PPTXPresentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    for slide_content in content:
        # Add a slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use layout with title and content
        
        # Set title
        title = slide.shapes.title
        title.text = slide_content['title']
        
        # Add points
        body = slide.shapes.placeholders[1]  # Second placeholder is for body
        tf = body.text_frame
        
        for point in slide_content['points']:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
        
        # If there's an image suggestion and it's a URL, try to add it
        if 'image' in slide_content and slide_content['image'].startswith('http'):
            try:
                response = requests.get(slide_content['image'])
                if response.status_code == 200:
                    img_stream = io.BytesIO(response.content)
                    
                    # Add watermark if specified
                    if watermark:
                        img_stream = add_watermark(img_stream, watermark)
                    
                    # Add image to slide
                    slide.shapes.add_picture(img_stream, Inches(8), Inches(2), height=Inches(4))
            except Exception as e:
                print(f"Error adding image: {str(e)}")
    
    # Save presentation to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes
