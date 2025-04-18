from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for, flash
import os
import json
from dotenv import load_dotenv
import logging
from datetime import datetime
from flask_wtf.csrf import CSRFProtect, generate_csrf, CSRFError
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, login_user, logout_user, login_required, current_user, UserMixin
from flask_cors import CORS
import tempfile
import requests
from PIL import Image
from io import BytesIO
import openai
from flask_dance.contrib.google import make_google_blueprint, google
from flask_dance.consumer.storage.sqla import OAuthConsumerMixin, SQLAlchemyStorage
from flask_dance.consumer import oauth_authorized
from sqlalchemy.orm.exc import NoResultFound
from utils.utils import check_user_limits, get_max_slides, add_watermark, generate_presentation_content, create_ppt
from pptx import Presentation
import time
from flask import send_from_directory
from functools import wraps

# Load environment variables
load_dotenv()

# Initialize extensions
db = SQLAlchemy()
login_manager = LoginManager()
csrf = CSRFProtect()

# Initialize Google Blueprint globally
google_bp = make_google_blueprint(
    client_id=os.getenv('GOOGLE_CLIENT_ID'),
    client_secret=os.getenv('GOOGLE_CLIENT_SECRET'),
    scope=['openid', 'https://www.googleapis.com/auth/userinfo.email', 'https://www.googleapis.com/auth/userinfo.profile'],
    redirect_to='index'  # Redirect to index page after login
)

# Configure static folder for Render deployment
if os.environ.get('RENDER'):
    # On Render, use a tmp directory that we have write access to
    STATIC_FOLDER = '/tmp/static'
    PRESENTATION_FOLDER = os.path.join(STATIC_FOLDER, 'presentations')
else:
    # Local development
    STATIC_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
    PRESENTATION_FOLDER = os.path.join(STATIC_FOLDER, 'presentations')

# Models
class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(120), unique=True, nullable=False)
    name = db.Column(db.String(120))
    google_id = db.Column(db.String(256), unique=True)

    def get_id(self):
        return str(self.id)

class OAuth(OAuthConsumerMixin, db.Model):
    __tablename__ = "flask_dance_oauth"
    provider_user_id = db.Column(db.String(256), unique=True)
    user_id = db.Column(db.Integer, db.ForeignKey(User.id))
    user = db.relationship(User)

    def set_token(self, token):
        app.logger.debug(f"Token type: {type(token)}")
        app.logger.debug(f"Token content: {token}")
        
        # Convert token to dictionary if it's not already
        if not isinstance(token, dict):
            try:
                if hasattr(token, 'token'):
                    token = token.token
                token = {
                    'access_token': getattr(token, 'access_token', None),
                    'token_type': getattr(token, 'token_type', None),
                    'scope': getattr(token, 'scope', []),
                    'expires_in': getattr(token, 'expires_in', None),
                    'expires_at': getattr(token, 'expires_at', None),
                    'id_token': getattr(token, 'id_token', None)
                }
            except Exception as e:
                app.logger.error(f"Error converting token to dictionary: {str(e)}")
                token = {}
        
        app.logger.debug(f"Token after conversion: {token}")
        
        try:
            # Clean up the token dictionary
            token_dict = {
                'access_token': str(token.get('access_token', '')),
                'token_type': str(token.get('token_type', '')),
                'scope': token.get('scope', []),
                'expires_in': int(token.get('expires_in', 0)),
                'expires_at': float(token.get('expires_at', 0)),
                'id_token': str(token.get('id_token', ''))
            }
            
            # Convert scope to string if it's a list
            if isinstance(token_dict['scope'], list):
                token_dict['scope'] = ' '.join(token_dict['scope'])
            
            app.logger.debug(f"Final token_dict: {token_dict}")
            json_token = json.dumps(token_dict)
            app.logger.debug(f"JSON token: {json_token}")
            self.token = json_token
            
        except Exception as e:
            app.logger.error(f"Error in token serialization: {str(e)}")
            self.token = json.dumps({})

    def get_token(self):
        try:
            app.logger.debug(f"Stored token: {self.token}")
            token_dict = json.loads(self.token)
            app.logger.debug(f"Loaded token dict: {token_dict}")
            # Convert scope back to list if it was stored as a string
            if isinstance(token_dict.get('scope'), str):
                token_dict['scope'] = token_dict['scope'].split()
            return token_dict
        except json.JSONDecodeError as e:
            app.logger.error(f"Failed to decode token JSON: {str(e)}")
            return {}
        except Exception as e:
            app.logger.error(f"Unexpected error in get_token: {str(e)}")
            return {}

class Presentation(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(120), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    user = db.relationship('User', backref=db.backref('presentations', lazy=True))
    file_path = db.Column(db.String(500))
    status = db.Column(db.String(50), default='pending')
    error_message = db.Column(db.Text)

def create_app():
    app = Flask(__name__, static_folder=STATIC_FOLDER)
    app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'your-secret-key')
    app.config['STATIC_FOLDER'] = STATIC_FOLDER
    app.config['PRESENTATION_FOLDER'] = PRESENTATION_FOLDER
    
    # Configure logging
    logging.basicConfig(level=logging.DEBUG)
    app.logger.setLevel(logging.DEBUG)

    # Configure database
    database_url = os.getenv('DATABASE_URL')
    if database_url and database_url.startswith('postgres://'):
        database_url = database_url.replace('postgres://', 'postgresql://', 1)

    app.config['SQLALCHEMY_DATABASE_URI'] = database_url or 'sqlite:///app.db'
    app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
    app.config['WTF_CSRF_ENABLED'] = True
    
    # Initialize extensions with app
    db.init_app(app)
    login_manager.init_app(app)
    csrf.init_app(app)
    
    # Register Google Blueprint
    app.register_blueprint(google_bp, url_prefix='/login')
    
    # Configure OAuth storage
    google_bp.storage = SQLAlchemyStorage(OAuth, db.session, user=current_user)
    
    # Create database tables
    with app.app_context():
        try:
            db.create_all()
            app.logger.info("Database tables created successfully")
        except Exception as e:
            app.logger.error(f"Error creating database tables: {str(e)}")
    
    # Ensure static folders exist
    os.makedirs(PRESENTATION_FOLDER, exist_ok=True)
    app.logger.info(f'Static folder configured at: {STATIC_FOLDER}')
    app.logger.info(f'Presentations folder configured at: {PRESENTATION_FOLDER}')
    
    return app

app = create_app()

# Custom CSRF error handler
@app.errorhandler(CSRFError)
def handle_csrf_error(e):
    app.logger.error(f"CSRF Error: {str(e)}")
    return jsonify(error="CSRF token validation failed"), 400

# Custom decorator for JSON CSRF validation
def csrf_protect_json():
    def decorator(view_function):
        @wraps(view_function)
        def wrapped_view(*args, **kwargs):
            # Check for CSRF token in X-CSRF-Token header
            token = request.headers.get('X-CSRF-Token')
            if not token:
                app.logger.error("Missing CSRF token in request headers")
                return jsonify(error="Missing CSRF token"), 400
            
            # Validate token using Flask-WTF's validate_csrf
            try:
                csrf.protect()
            except CSRFError as e:
                app.logger.error(f"CSRF validation failed: {str(e)}")
                return jsonify(error="Invalid CSRF token"), 400
                
            return view_function(*args, **kwargs)
        return wrapped_view
    return decorator

# Add CSRF token to all responses
@app.after_request
def add_csrf_token(response):
    if 'text/html' in response.headers.get('Content-Type', ''):
        token = generate_csrf()  # Using Flask-WTF's generate_csrf
        app.logger.debug(f"Generated new CSRF token")
        response.set_cookie('csrf_token', token, secure=True, httponly=True, samesite='Strict')
    return response

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

@app.before_request
def before_request():
    if not request.is_secure and app.env != 'development':
        url = request.url.replace('http://', 'https://', 1)
        return redirect(url, code=301)

@oauth_authorized.connect_via(google_bp)
def google_logged_in(blueprint, token):
    try:
        if not token:
            app.logger.error("Failed to get OAuth token")
            return False

        resp = blueprint.session.get('/oauth2/v2/userinfo')
        if not resp.ok:
            app.logger.error(f"Failed to get user info: {resp.text}")
            return False

        google_info = resp.json()
        email = google_info['email']
        
        # Find or create user
        user = User.query.filter_by(email=email).first()
        if not user:
            user = User(
                email=email,
                name=google_info.get('name', email.split('@')[0]),
                google_id=google_info['id']
            )
            db.session.add(user)
            
            # Create OAuth entry
            oauth = OAuth(
                provider=blueprint.name,
                provider_user_id=google_info['id'],
                token=token,
                user=user
            )
            db.session.add(oauth)
            
            try:
                db.session.commit()
            except Exception as e:
                app.logger.error(f"Error saving user: {str(e)}")
                db.session.rollback()
                return False
        
        # Log in the user
        login_user(user)
        app.logger.info(f"Successfully logged in user {email}")
        
        # Return False to disable Flask-Dance's default behavior
        return False
        
    except Exception as e:
        app.logger.error(f"Error in OAuth callback: {str(e)}")
        return False

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('index'))

@app.route('/login')
def login():
    return redirect(url_for('google.login'))

@app.route('/register', methods=['GET', 'POST'])
def register():
    return jsonify({
        'status': 'error',
        'message': 'Registration is not available. Please use Google to sign in.'
    }), 400

@app.route('/verify-email/<token>')
def verify_email(token):
    return jsonify({
        'status': 'error',
        'message': 'Email verification is not available. Please use Google to sign in.'
    }), 400

@app.route('/resend-verification', methods=['POST'])
def resend_verification():
    return jsonify({
        'status': 'error',
        'message': 'Email verification is not available. Please use Google to sign in.'
    }), 400

@app.route('/static/presentations/<path:filename>')
def serve_presentation(filename):
    app.logger.info(f'Serving presentation file: {filename}')
    try:
        return send_from_directory(app.config['PRESENTATION_FOLDER'], filename, as_attachment=True)
    except Exception as e:
        app.logger.error(f'Error serving file {filename}: {str(e)}')
        return jsonify({'error': 'File not found'}), 404

@app.route('/generate', methods=['POST'])
@login_required
@csrf_protect_json()
def generate_presentation():
    app.logger.info(f"Generate presentation request received from user: {current_user.email}")
    app.logger.debug(f"Request headers: {dict(request.headers)}")
    app.logger.debug(f"Request method: {request.method}")
    app.logger.debug(f"Content-Type: {request.content_type}")
    
    try:
        # Handle both JSON and form data
        if request.is_json:
            app.logger.debug("Processing JSON request")
            data = request.get_json()
        else:
            app.logger.debug("Processing form data request")
            data = {
                'topic': request.form.get('topic'),
                'num_slides': request.form.get('slides'),
                'theme': request.form.get('theme')
            }
        
        app.logger.debug(f"Received data: {data}")
        
        if not data:
            app.logger.error("No data received in request")
            return jsonify({"success": False, "error": "No data received"}), 400
            
        topic = data.get('topic')
        num_slides = int(data.get('num_slides', data.get('slides', 5)))
        theme = data.get('theme', 'professional')
        
        if not topic:
            app.logger.error("Missing topic in request")
            return jsonify({"success": False, "error": "Topic is required"}), 400
            
        app.logger.info(f"Generating presentation: {topic}, {num_slides} slides, {theme} theme")
        
        # Create presentations directory if it doesn't exist
        if not os.path.exists(app.config['PRESENTATION_FOLDER']):
            app.logger.info(f"Creating presentations directory: {app.config['PRESENTATION_FOLDER']}")
            os.makedirs(app.config['PRESENTATION_FOLDER'])
        
        # Generate unique filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"presentation_{timestamp}_{current_user.id}.pptx"
        filepath = os.path.join(app.config['PRESENTATION_FOLDER'], filename)
        
        app.logger.info(f"Saving presentation to: {filepath}")
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic
        subtitle.text = f"Generated for {current_user.name}"
        
        # Save the presentation
        try:
            prs.save(filepath)
            app.logger.info("Presentation saved successfully")
            
            # Generate download URL
            download_url = url_for('serve_presentation', filename=filename)
            app.logger.info(f"Download URL: {download_url}")
            
            return jsonify({
                "success": True,
                "message": "Presentation generated successfully",
                "download_url": download_url
            })
            
        except Exception as e:
            app.logger.error(f"Error saving presentation: {str(e)}")
            return jsonify({"success": False, "error": f"Error saving presentation: {str(e)}"}), 500
            
    except Exception as e:
        app.logger.error(f"Error generating presentation: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/pricing')
def pricing():
    return render_template('pricing.html', plans=['free', 'pro', 'business'])

@app.route('/subscribe', methods=['POST'])
@login_required
def subscribe():
    try:
        data = request.get_json()
        plan = data.get('plan')
        
        if not plan or plan not in ['free', 'pro', 'business']:
            return jsonify({
                'status': 'error',
                'message': 'Invalid subscription plan'
            }), 400
            
        if plan == 'free':
            current_user.subscription_type = 'free'
            db.session.commit()
            return jsonify({
                'status': 'success',
                'message': 'Successfully subscribed to Free plan'
            })
            
        # For paid plans, initialize payment with Paystack
        amount = {
            'pro': 1000,  # ₦1,000 per month
            'business': 5000  # ₦5,000 per month
        }.get(plan)
        
        url = "https://api.paystack.co/transaction/initialize"
        headers = {
            "Authorization": f"Bearer {app.config['PAYSTACK_SECRET_KEY']}",
            "Content-Type": "application/json"
        }
        data = {
            "email": current_user.email,
            "amount": amount * 100,  # Amount in kobo
            "callback_url": url_for('payment_callback', _external=True),
            "metadata": {
                "plan": plan,
                "user_id": current_user.id
            }
        }
        
        response = requests.post(url, headers=headers, json=data)
        if response.ok:
            data = response.json()
            if data.get('status'):
                return jsonify({
                    'status': 'success',
                    'authorization_url': data['data']['authorization_url']
                })
                
        return jsonify({
            'status': 'error',
            'message': 'Failed to initialize payment'
        }), 500
        
    except Exception as e:
        app.logger.error(f"Subscription error: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'An error occurred while processing your subscription'
        }), 500

@app.route('/payment/callback')
@login_required
def payment_callback():
    reference = request.args.get('reference')
    if not reference:
        flash('No reference provided')
        return redirect(url_for('pricing'))
        
    try:
        # Verify the transaction
        response = verify_paystack_transaction(reference)
        
        if response.get('status') and response['data']['status'] == 'success':
            metadata = response['data']['metadata']
            plan = metadata.get('plan')
            
            if plan in ['pro', 'business']:
                current_user.subscription_type = plan
                db.session.commit()
                flash(f'Successfully subscribed to {plan.title()} plan!')
            else:
                flash('Invalid subscription plan')
        else:
            flash('Payment verification failed')
            
    except Exception as e:
        app.logger.error(f"Payment callback error: {str(e)}")
        flash('An error occurred while processing your payment')
        
    return redirect(url_for('pricing'))

@app.route('/download/<filename>')
@login_required
def download_presentation(filename):
    try:
        # Get the presentation from the database
        presentation = Presentation.query.filter(
            Presentation.user_id == current_user.id,
            Presentation.file_path.like(f"%{filename}")
        ).first()
        
        if not presentation:
            return jsonify({
                'status': 'error',
                'message': 'Presentation not found'
            }), 404
            
        if not os.path.exists(presentation.file_path):
            return jsonify({
                'status': 'error',
                'message': 'Presentation file not found'
            }), 404
            
        return send_file(
            presentation.file_path,
            as_attachment=True,
            download_name='presentation.pptx'
        )
        
    except Exception as e:
        app.logger.error(f"Error downloading presentation: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Failed to download presentation'
        }), 500

@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    return jsonify({
        'status': 'error',
        'message': 'Password reset is not available. Please use Google to sign in.'
    }), 400

@app.route('/reset-password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    return jsonify({
        'status': 'error',
        'message': 'Password reset is not available. Please use Google to sign in.'
    }), 400

if __name__ == '__main__':
    app.run(debug=True)
